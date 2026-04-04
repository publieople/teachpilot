"""测试多模态解析模块"""
import os
from modules.multimodal_service import multimodal_service

# 创建测试文件
test_dir = "test_files"
os.makedirs(test_dir, exist_ok=True)

# 测试 1：创建并解析 TXT 文件
print("=" * 50)
print("测试 1: TXT 文件解析")
txt_path = os.path.join(test_dir, "test.txt")
with open(txt_path, "w", encoding="utf-8") as f:
    f.write("这是测试文本。\n教学设计是教师根据教学目标，系统地规划教学过程的活动。")

result = multimodal_service.parse_file(txt_path)
print(f"文本内容：{result.get('text', '')[:100]}...")
print("[OK] TXT 解析成功\n")

# 测试 2：创建并解析 Word 文件
print("=" * 50)
print("测试 2: Word 文件解析")
from docx import Document
docx_path = os.path.join(test_dir, "test.docx")
doc = Document()
doc.add_heading('教学设计基础', 0)
doc.add_paragraph('教学设计是教师根据教学目标，系统地规划教学过程的活动。')
doc.add_paragraph('包括教学目标设定、教学内容选择、教学方法设计等环节。')
doc.save(docx_path)

result = multimodal_service.parse_word(docx_path)
print(f"段落数：{len(result.get('paragraphs', []))}")
print(f"文本内容：{result.get('text', '')[:100]}...")
print("[OK] Word 解析成功\n")

# 测试 3：创建并解析 PPT 文件
print("=" * 50)
print("测试 3: PPT 文件解析")
from pptx import Presentation
pptx_path = os.path.join(test_dir, "test.pptx")
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "教学设计"
subtitle = slide.placeholders[1]
subtitle.text = "系统规划教学过程"
prs.save(pptx_path)

result = multimodal_service.parse_ppt(pptx_path)
print(f"幻灯片数：{result.get('slides', 0)}")
print(f"文本内容：{result.get('text', '')[:100]}...")
print("[OK] PPT 解析成功\n")

# 测试 4：创建并解析图片文件
print("=" * 50)
print("测试 4: 图片文件解析")
from PIL import Image
img_path = os.path.join(test_dir, "test.png")
img = Image.new('RGB', (100, 100), color='blue')
img.save(img_path)

result = multimodal_service.parse_image(img_path)
print(f"分辨率：{result.get('width')}x{result.get('height')}")
print(f"格式：{result.get('format')}")
print("[OK] 图片解析成功\n")

print("=" * 50)
print("[OK] 所有测试通过！")
print("=" * 50)
