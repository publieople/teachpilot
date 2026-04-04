"""
多模态文件解析服务
支持：PDF, Word, PPT, 图片，视频
"""

import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import cv2
import base64
from io import BytesIO
from typing import Dict, List, Any, Optional
from pathlib import Path


class MultimodalService:
    """多模态文件解析服务"""
    
    def __init__(self):
        """初始化服务"""
        self.supported_extensions = {
            'pdf': ['.pdf'],
            'word': ['.docx', '.doc'],
            'ppt': ['.pptx', '.ppt'],
            'image': ['.jpg', '.jpeg', '.png', '.gif', '.webp'],
            'video': ['.mp4', '.avi', '.mov', '.mkv'],
            'text': ['.txt', '.md', '.markdown']
        }
    
    def parse_file(self, file_path: str, file_type: Optional[str] = None) -> Dict[str, Any]:
        """
        解析文件
        
        Args:
            file_path: 文件路径
            file_type: 文件类型（可选，自动检测）
            
        Returns:
            解析结果字典
        """
        if not file_type:
            file_type = self._detect_file_type(file_path)
        
        if file_type == 'pdf':
            return self.parse_pdf(file_path)
        elif file_type in ['word', 'docx']:
            return self.parse_word(file_path)
        elif file_type in ['ppt', 'pptx']:
            return self.parse_ppt(file_path)
        elif file_type == 'image':
            return self.parse_image(file_path)
        elif file_type == 'video':
            return self.parse_video(file_path)
        elif file_type == 'text':
            return self.parse_text(file_path)
        else:
            raise ValueError(f"不支持的文件类型：{file_type}")
    
    def _detect_file_type(self, file_path: str) -> str:
        """检测文件类型"""
        ext = Path(file_path).suffix.lower()
        
        for file_type, extensions in self.supported_extensions.items():
            if ext in extensions:
                return file_type
        
        raise ValueError(f"未知的文件类型：{ext}")
    
    def parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        解析 PDF 文件
        
        Returns:
            {
                'text': str,  # 全部文本
                'pages': int,  # 页数
                'metadata': dict,  # 元数据
                'images': List[dict]  # 提取的图片
            }
        """
        doc = fitz.open(file_path)
        
        result = {
            'text': '',
            'pages': len(doc),
            'metadata': {
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'subject': doc.metadata.get('subject', ''),
            },
            'images': [],
            'content_by_page': []
        }
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # 提取文本
            page_text = page.get_text()
            result['text'] += page_text + "\n"
            result['content_by_page'].append({
                'page': page_num + 1,
                'text': page_text
            })
            
            # 提取图片
            image_list = page.get_images(full=True)
            for img_index, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    result['images'].append({
                        'page': page_num + 1,
                        'index': img_index,
                        'format': base_image["ext"],
                        'size': len(image_bytes)
                    })
                except:
                    pass
        
        doc.close()
        return result
    
    def parse_word(self, file_path: str) -> Dict[str, Any]:
        """
        解析 Word 文件
        
        Returns:
            {
                'text': str,  # 全部文本
                'paragraphs': List[str],  # 段落列表
                'tables': List[List[List[str]]],  # 表格
                'metadata': dict
            }
        """
        doc = Document(file_path)
        
        result = {
            'text': '',
            'paragraphs': [],
            'tables': [],
            'metadata': {
                'author': doc.core_properties.author if hasattr(doc.core_properties, 'author') else '',
                'title': doc.core_properties.title if hasattr(doc.core_properties, 'title') else '',
            }
        }
        
        # 提取段落
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                result['paragraphs'].append(text)
                result['text'] += text + "\n"
        
        # 提取表格
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            if table_data:
                result['tables'].append(table_data)
        
        return result
    
    def parse_ppt(self, file_path: str) -> Dict[str, Any]:
        """
        解析 PPT 文件
        
        Returns:
            {
                'text': str,  # 全部文本
                'slides': int,  # 幻灯片数
                'content_by_slide': List[dict],  # 每页内容
                'metadata': dict
            }
        """
        prs = Presentation(file_path)
        
        result = {
            'text': '',
            'slides': len(prs.slides),
            'content_by_slide': [],
            'metadata': {}
        }
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = ''
            shapes_info = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"
                    shapes_info.append({
                        'type': shape.shape_type,
                        'text': shape.text
                    })
            
            result['text'] += slide_text + "\n"
            result['content_by_slide'].append({
                'slide': slide_num + 1,
                'text': slide_text,
                'shapes': len(shapes_info)
            })
        
        return result
    
    def parse_text(self, file_path: str) -> Dict[str, Any]:
        """
        解析文本文件（TXT/MD）
        
        Returns:
            {
                'text': str,  # 全部文本
                'lines': int,  # 行数
                'size': int  # 文件大小
            }
        """
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        return {
            'text': content,
            'lines': len(content.splitlines()),
            'size': os.path.getsize(file_path)
        }
    
    def parse_image(self, file_path: str) -> Dict[str, Any]:
        """
        解析图片文件
        
        Returns:
            {
                'width': int,
                'height': int,
                'format': str,
                'size': int,
                'metadata': dict
            }
        """
        img = Image.open(file_path)
        file_size = os.path.getsize(file_path)
        
        return {
            'width': img.width,
            'height': img.height,
            'format': img.format,
            'mode': img.mode,
            'size': file_size,
            'metadata': {
                'info': dict(img.info) if img.info else {}
            }
        }
    
    def parse_video(self, file_path: str, num_frames: int = 10) -> Dict[str, Any]:
        """
        解析视频文件（提取关键帧）
        
        Args:
            file_path: 视频文件路径
            num_frames: 提取的关键帧数量
            
        Returns:
            {
                'duration': float,  # 时长（秒）
                'fps': float,  # 帧率
                'width': int,
                'height': int,
                'total_frames': int,
                'keyframes': List[dict]  # 关键帧信息
            }
        """
        cap = cv2.VideoCapture(file_path)
        
        # 获取视频信息
        fps = cap.get(cv2.CAP_PROP_FPS)
        frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        duration = frame_count / fps if fps > 0 else 0
        
        result = {
            'duration': round(duration, 2),
            'fps': round(fps, 2),
            'width': width,
            'height': height,
            'total_frames': frame_count,
            'keyframes': []
        }
        
        # 提取关键帧（均匀采样）
        if frame_count > 0:
            frame_indices = self._generate_frame_indices(frame_count, num_frames)
            
            for idx, frame_idx in enumerate(frame_indices):
                cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
                ret, frame = cap.read()
                
                if ret:
                    # 保存关键帧信息
                    result['keyframes'].append({
                        'index': idx,
                        'frame_number': frame_idx,
                        'timestamp': round(frame_idx / fps, 2) if fps > 0 else 0,
                        'shape': frame.shape
                    })
        
        cap.release()
        return result
    
    def _generate_frame_indices(self, total_frames: int, num_frames: int) -> List[int]:
        """生成均匀采样的帧索引"""
        if total_frames <= num_frames:
            return list(range(total_frames))
        
        step = total_frames // num_frames
        return [i * step for i in range(num_frames)]
    
    def get_file_summary(self, file_path: str) -> str:
        """
        获取文件摘要（用于 RAG）
        
        Returns:
            文件内容摘要文本
        """
        result = self.parse_file(file_path)
        
        # 根据文件类型生成摘要
        if 'text' in result:
            # 返回前 1000 字符
            return result['text'][:1000] if len(result['text']) > 1000 else result['text']
        elif 'duration' in result:
            # 视频文件
            return f"视频文件，时长 {result['duration']} 秒，分辨率 {result['width']}x{result['height']}，共 {result['total_frames']} 帧"
        elif 'width' in result:
            # 图片文件
            return f"图片文件，分辨率 {result['width']}x{result['height']}，格式 {result['format']}"
        
        return ""


# 全局实例
multimodal_service = MultimodalService()
