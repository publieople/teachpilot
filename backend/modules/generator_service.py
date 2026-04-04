"""
课件生成服务
支持：PPT 生成、Word 教案生成、动画/互动内容
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt, Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from typing import Dict, List, Any, Optional
from datetime import datetime


class GeneratorService:
    """课件生成服务"""
    
    def __init__(self):
        """初始化服务"""
        self.output_dir = "outputs"
        os.makedirs(self.output_dir, exist_ok=True)
    
    # ==================== PPT 生成 ====================
    
    def generate_ppt(
        self,
        title: str,
        slides: List[Dict[str, Any]],
        subtitle: Optional[str] = None,
        output_filename: Optional[str] = None
    ) -> str:
        """
        生成 PPT 课件
        
        Args:
            title: 课程标题
            slides: 幻灯片列表，每项包含：
                - title: 标题
                - content: 内容（列表或字符串）
                - notes: 备注（可选）
            subtitle: 副标题（可选）
            output_filename: 输出文件名（可选）
            
        Returns:
            生成的 PPT 文件路径
        """
        prs = Presentation()
        
        # 1. 封面页
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle
        else:
            subtitle_shape.text = f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # 2. 目录页
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "目录"
        tf = content_shape.text_frame
        tf.clear()
        
        for i, slide_data in enumerate(slides, 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {slide_data.get('title', '无标题')}"
            p.level = 0
            p.font.size = Pt(18)
        
        # 3. 内容页
        for slide_data in slides:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            # 设置标题
            title_shape.text = slide_data.get('title', '')
            
            # 设置内容
            content = slide_data.get('content', '')
            tf = content_shape.text_frame
            tf.clear()
            
            if isinstance(content, list):
                for item in content:
                    p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.level = 0
                    p.font.size = Pt(16)
                    p.space_after = Pt(10)
            else:
                p = tf.add_paragraph()
                p.text = str(content)
                p.font.size = Pt(16)
            
            # 添加备注
            if 'notes' in slide_data:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data['notes']
        
        # 4. 总结页
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "总结"
        tf = content_shape.text_frame
        tf.clear()
        
        p = tf.add_paragraph()
        p.text = "本课程主要内容回顾："
        p.font.size = Pt(18)
        p.font.bold = True
        
        for i, slide_data in enumerate(slides, 1):
            p = tf.add_paragraph()
            p.text = f"• {slide_data.get('title', '无标题')}"
            p.level = 0
            p.font.size = Pt(16)
        
        # 保存文件
        if not output_filename:
            output_filename = f"ppt_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        output_path = os.path.join(self.output_dir, output_filename)
        prs.save(output_path)
        
        return output_path
    
    # ==================== Word 教案生成 ====================
    
    def generate_word_lesson_plan(
        self,
        title: str,
        teaching_objectives: str,
        teaching_content: str,
        teaching_methods: str,
        teaching_process: List[Dict[str, str]],
        classroom_activities: str,
        homework: str,
        output_filename: Optional[str] = None
    ) -> str:
        """
        生成 Word 教案
        
        Args:
            title: 课程标题
            teaching_objectives: 教学目标
            teaching_content: 教学内容
            teaching_methods: 教学方法
            teaching_process: 教学过程（列表，每项包含 stage 和 content）
            classroom_activities: 课堂活动
            homework: 课后作业
            output_filename: 输出文件名
            
        Returns:
            生成的 Word 文件路径
        """
        doc = Document()
        
        # 标题
        heading = doc.add_heading(title, 0)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 基本信息
        doc.add_heading('一、基本信息', level=1)
        p = doc.add_paragraph()
        p.add_run(f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        
        # 教学目标
        doc.add_heading('二、教学目标', level=1)
        doc.add_paragraph(teaching_objectives)
        
        # 教学内容
        doc.add_heading('三、教学内容', level=1)
        doc.add_paragraph(teaching_content)
        
        # 教学方法
        doc.add_heading('四、教学方法', level=1)
        doc.add_paragraph(teaching_methods)
        
        # 教学过程
        doc.add_heading('五、教学过程', level=1)
        for i, stage in enumerate(teaching_process, 1):
            p = doc.add_paragraph()
            p.add_run(f"{i}. {stage.get('stage', '环节')}").bold = True
            p.add_run(f"\n   {stage.get('content', '')}")
        
        # 课堂活动
        doc.add_heading('六、课堂活动', level=1)
        doc.add_paragraph(classroom_activities)
        
        # 课后作业
        doc.add_heading('七、课后作业', level=1)
        doc.add_paragraph(homework)
        
        # 教学反思（预留）
        doc.add_heading('八、教学反思', level=1)
        doc.add_paragraph("（课后填写）")
        
        # 保存文件
        if not output_filename:
            output_filename = f"lesson_plan_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        
        output_path = os.path.join(self.output_dir, output_filename)
        doc.save(output_path)
        
        return output_path
    
    # ==================== 动画/互动内容生成 ====================
    
    def generate_html_animation(
        self,
        title: str,
        content: str,
        animation_type: str = "fade"
    ) -> str:
        """
        生成简单的 HTML 动画
        
        Args:
            title: 标题
            content: 内容
            animation_type: 动画类型（fade, slide, bounce）
            
        Returns:
            HTML 文件路径
        """
        animations = {
            'fade': '''
                @keyframes fadeIn {
                    from { opacity: 0; }
                    to { opacity: 1; }
                }
                .animated { animation: fadeIn 2s ease-in; }
            ''',
            'slide': '''
                @keyframes slideIn {
                    from { transform: translateX(-100%); }
                    to { transform: translateX(0); }
                }
                .animated { animation: slideIn 1.5s ease-out; }
            ''',
            'bounce': '''
                @keyframes bounceIn {
                    0%, 20%, 40%, 60%, 80%, 100% { transition-timing-function: cubic-bezier(0.215, 0.610, 0.355, 1.000); }
                    0% { opacity: 0; transform: scale3d(.3, .3, .3); }
                    20% { transform: scale3d(1.1, 1.1, 1.1); }
                    40% { transform: scale3d(.9, .9, .9); }
                    60% { opacity: 1; transform: scale3d(1.03, 1.03, 1.03); }
                    80% { transform: scale3d(.97, .97, .97); }
                    100% { opacity: 1; transform: scale3d(1, 1, 1); }
                }
                .animated { animation: bounceIn 2s; }
            '''
        }
        
        animation_css = animations.get(animation_type, animations['fade'])
        
        html_content = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{
            font-family: 'Arial', sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            display: flex;
            justify-content: center;
            align-items: center;
            margin: 0;
        }}
        .container {{
            background: white;
            padding: 40px;
            border-radius: 20px;
            box-shadow: 0 20px 60px rgba(0,0,0,0.3);
            max-width: 800px;
            width: 90%;
        }}
        h1 {{
            color: #667eea;
            text-align: center;
            margin-bottom: 30px;
        }}
        .content {{
            font-size: 18px;
            line-height: 1.8;
            color: #333;
        }}
        {animation_css}
    </style>
</head>
<body>
    <div class="container">
        <h1 class="animated">{title}</h1>
        <div class="content animated">
            {content.replace(chr(10), '<br>')}
        </div>
    </div>
</body>
</html>'''
        
        output_filename = f"animation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        output_path = os.path.join(self.output_dir, output_filename)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return output_path
    
    def generate_quiz_game(
        self,
        title: str,
        questions: List[Dict[str, Any]]
    ) -> str:
        """
        生成互动问答游戏（HTML）
        
        Args:
            title: 游戏标题
            questions: 问题列表，每项包含：
                - question: 问题
                - options: 选项列表
                - answer: 正确答案索引
                
        Returns:
            HTML 文件路径
        """
        questions_json = str(questions)
        
        html_content = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title} - 互动问答</title>
    <style>
        body {{
            font-family: 'Arial', sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            display: flex;
            justify-content: center;
            align-items: center;
            margin: 0;
        }}
        .container {{
            background: white;
            padding: 40px;
            border-radius: 20px;
            box-shadow: 0 20px 60px rgba(0,0,0,0.3);
            max-width: 600px;
            width: 90%;
        }}
        h1 {{
            color: #667eea;
            text-align: center;
        }}
        .question {{
            font-size: 20px;
            margin: 20px 0;
            color: #333;
        }}
        .options {{
            display: flex;
            flex-direction: column;
            gap: 10px;
        }}
        .option {{
            padding: 15px 20px;
            background: #f0f0f0;
            border: 2px solid #ddd;
            border-radius: 10px;
            cursor: pointer;
            transition: all 0.3s;
        }}
        .option:hover {{
            background: #667eea;
            color: white;
            border-color: #667eea;
        }}
        .option.correct {{
            background: #4CAF50;
            color: white;
            border-color: #4CAF50;
        }}
        .option.wrong {{
            background: #f44336;
            color: white;
            border-color: #f44336;
        }}
        .score {{
            text-align: center;
            font-size: 24px;
            color: #667eea;
            margin-top: 20px;
        }}
        .hidden {{
            display: none;
        }}
    </style>
</head>
<body>
    <div class="container">
        <h1>{title}</h1>
        <div id="quiz">
            <div class="question" id="question"></div>
            <div class="options" id="options"></div>
        </div>
        <div id="result" class="hidden">
            <h2>游戏结束!</h2>
            <div class="score" id="score"></div>
        </div>
    </div>
    
    <script>
        const questions = {questions_json};
        let currentQuestion = 0;
        let score = 0;
        
        function loadQuestion() {{
            if (currentQuestion >= questions.length) {{
                showResult();
                return;
            }}
            
            const q = questions[currentQuestion];
            document.getElementById('question').textContent = `第${{currentQuestion + 1}}题：${{q.question}}`;
            
            const optionsDiv = document.getElementById('options');
            optionsDiv.innerHTML = '';
            
            q.options.forEach((opt, index) => {{
                const div = document.createElement('div');
                div.className = 'option';
                div.textContent = opt;
                div.onclick = () => checkAnswer(index, div);
                optionsDiv.appendChild(div);
            }});
        }}
        
        function checkAnswer(selectedIndex, element) {{
            const q = questions[currentQuestion];
            const options = document.querySelectorAll('.option');
            
            options.forEach(opt => opt.style.pointerEvents = 'none');
            
            if (selectedIndex === q.answer) {{
                element.classList.add('correct');
                score++;
            }} else {{
                element.classList.add('wrong');
                options[q.answer].classList.add('correct');
            }}
            
            setTimeout(() => {{
                currentQuestion++;
                loadQuestion();
            }}, 1500);
        }}
        
        function showResult() {{
            document.getElementById('quiz').classList.add('hidden');
            document.getElementById('result').classList.remove('hidden');
            document.getElementById('score').textContent = `得分：${{score}} / ${{questions.length}}`;
        }}
        
        loadQuestion();
    </script>
</body>
</html>'''
        
        output_filename = f"quiz_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        output_path = os.path.join(self.output_dir, output_filename)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return output_path


# 全局实例
generator_service = GeneratorService()
